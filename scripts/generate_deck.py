from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io

# ── Brand palette ────────────────────────────────────────────────────────────
BG       = RGBColor(0x11, 0x11, 0x13)
SURFACE  = RGBColor(0x18, 0x18, 0x1b)
BORDER   = RGBColor(0x2a, 0x2a, 0x2e)
GREEN    = RGBColor(0x5d, 0x99, 0x80)
SAND     = RGBColor(0xc4, 0xb4, 0x9a)
TEXT     = RGBColor(0xe4, 0xe4, 0xe7)
TEXT_SEC = RGBColor(0x8b, 0x8b, 0x92)

A   = '/Users/beckettdevoe/Desktop/projects/landex/landing_page/public/assets/'
OUT = '/Users/beckettdevoe/Desktop/Landex_NVIDIA_Insight_2026.pptx'

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]   # completely blank


# ── Image helpers ─────────────────────────────────────────────────────────────

def prep_img(path, tw_in, th_in, darken=0.0):
    """Crop image to target aspect ratio, optionally darken, return BytesIO JPEG."""
    raw = Image.open(path)
    # Composite RGBA onto dark surface background
    if raw.mode == 'RGBA':
        bg = Image.new('RGB', raw.size, (0x18, 0x18, 0x1b))
        bg.paste(raw, mask=raw.split()[3])
        img = bg
    else:
        img = raw.convert('RGB')

    w, h = img.size
    tr, sr = tw_in / th_in, w / h
    if sr > tr:
        nw = int(h * tr)
        img = img.crop(((w - nw) // 2, 0, (w - nw) // 2 + nw, h))
    else:
        nh = int(w / tr)
        img = img.crop((0, (h - nh) // 2, w, (h - nh) // 2 + nh))

    # Resize to 150 dpi equivalent
    img = img.resize((int(tw_in * 150), int(th_in * 150)), Image.LANCZOS)

    if darken:
        black = Image.new('RGB', img.size, (0, 0, 0))
        img = Image.blend(img, black, darken)

    buf = io.BytesIO()
    img.save(buf, format='JPEG', quality=93)
    buf.seek(0)
    return buf


def sq_crop(path):
    """Center-crop to square, preserve transparency, return BytesIO PNG."""
    img = Image.open(path).convert('RGBA')
    w, h = img.size
    s = min(w, h)
    img = img.crop(((w - s) // 2, (h - s) // 2, (w - s) // 2 + s, (h - s) // 2 + s))
    img = img.resize((600, 600), Image.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf


# ── Shape / text helpers ──────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(blank)


def fill_rect(sl, x, y, w, h, color=BG, border=None):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s


def txt(sl, t, x, y, w, h, sz=11, bold=False, color=TEXT,
        align=PP_ALIGN.LEFT, italic=False, font='Syne'):
    txb = sl.shapes.add_textbox(x, y, w, h)
    tf  = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = t
    r.font.size    = Pt(sz)
    r.font.bold    = bold
    r.font.italic  = italic
    r.font.color.rgb = color
    r.font.name    = font
    return txb


def eyebrow(sl, t, x=Inches(0.65), y=Inches(0.38)):
    txt(sl, t.upper(), x, y, Inches(10), Inches(0.32),
        sz=7.5, color=GREEN, bold=True)


def hrule(sl, y, x=Inches(0.65), w=Inches(12.03)):
    fill_rect(sl, x, y, w, Pt(1), color=BORDER)


def green_tick(sl, x, y):
    fill_rect(sl, x, y, Pt(3), Inches(0.24), color=GREEN)


def add_logo(sl, x, y, width=Inches(1.4)):
    sl.shapes.add_picture(A + 'logo.png', x, y,
                          width=width, height=width * (150 / 500))


def logo_small(sl):
    add_logo(sl, W - Inches(1.75), Inches(0.22), width=Inches(1.4))


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — COVER
# ═══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()

# Left dark panel
fill_rect(s1, 0, 0, Inches(7.9), H)

# Right: darkened construction photo (full slide height)
right_photo = prep_img(
    A + 'frames-for-your-heart-VoI2jd75M6Q-unsplash.jpg',
    5.43, 7.5, darken=0.48
)
s1.shapes.add_picture(right_photo, Inches(7.9), 0,
                      width=Inches(5.43), height=H)

# Thin green vertical accent at seam
fill_rect(s1, Inches(7.88), 0, Pt(2), H, color=GREEN)

# NVIDIA eyebrow
txt(s1, 'NVIDIA INSIGHT PROGRAM  ·  2026',
    Inches(0.65), Inches(0.62), Inches(7), Inches(0.32),
    sz=8, color=GREEN, bold=True)

# Logo — large
add_logo(s1, Inches(0.65), Inches(1.12), width=Inches(2.7))

# Green accent rule under logo
fill_rect(s1, Inches(0.65), Inches(2.1), Inches(1.0), Pt(2.5), color=GREEN)

# Main title
txt(s1,
    'Turn an ordinary walk video\ninto a verified install record.',
    Inches(0.65), Inches(2.28), Inches(7.05), Inches(2.0),
    sz=34, bold=True, color=TEXT)

# Tagline
txt(s1,
    'Computer vision for construction verification — auto-tagging installed '
    'components against the approved submittal and writing each one back to '
    'the 3D model as a timestamped, attributable install record.',
    Inches(0.65), Inches(4.45), Inches(6.95), Inches(1.6),
    sz=12, color=TEXT_SEC)

# Contact footer
txt(s1, 'allen@landexsystems.com  ·  landexsystems.com',
    Inches(0.65), Inches(6.92), Inches(7), Inches(0.36),
    sz=9, color=TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
fill_rect(s2, 0, 0, W, H)
logo_small(s2)
eyebrow(s2, 'The Problem')

txt(s2, 'The walkdown already happens.\nThe record does not.',
    Inches(0.65), Inches(0.78), Inches(8.8), Inches(1.28),
    sz=30, bold=True, color=TEXT)

# Right photo strip — darkened construction photo
strip = prep_img(
    A + 'frames-for-your-heart-VoI2jd75M6Q-unsplash.jpg',
    4.45, 5.22, darken=0.52
)
s2.shapes.add_picture(strip, Inches(8.55), Inches(2.28),
                      width=Inches(4.45), height=Inches(5.22))
# Thin green accent line on left edge of photo strip
fill_rect(s2, Inches(8.53), Inches(2.28), Pt(2), Inches(5.22), color=GREEN)

hrule(s2, Inches(2.22))

problems = [
    ('The record gets rebuilt by hand.',
     'Engineers walk the area, check installed components against the approved submittal, '
     'return with notes — then retype them into the BIM element by element. '
     'The ground truth exists. It just is not captured.'),
    ('Evidence has no chain of custody.',
     'When a pay application is disputed or closeout stalls, proof is scattered across '
     'inboxes. No element-level audit trail, no timestamped evidence, no submittal linkage.'),
    ('Context never reaches the element.',
     'Verifying one valve means holding the submittal, spec, and change orders in your head '
     'simultaneously. No system assembles this at the moment of inspection.'),
]

y = Inches(2.42)
for title, body in problems:
    green_tick(s2, Inches(0.65), y + Inches(0.08))
    txt(s2, title, Inches(0.9), y, Inches(7.35), Inches(0.33),
        sz=12, bold=True, color=SAND)
    txt(s2, body, Inches(0.9), y + Inches(0.32), Inches(7.35), Inches(0.72),
        sz=10.5, color=TEXT_SEC)
    y += Inches(1.2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — TECHNICAL APPROACH
# ═══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
fill_rect(s3, 0, 0, W, H)

# Right panel — SURFACE card behind screenshot
fill_rect(s3, Inches(6.85), 0, Inches(6.48), H, color=SURFACE, border=BORDER)

# Product screenshot: video_annotator — 5.88" wide, ratio 2818/1890 = 1.4915
vid_h = 5.88 / (2818 / 1890)   # = 3.942"
vid_img = prep_img(A + 'video_annotator.png', 5.88, vid_h)
s3.shapes.add_picture(vid_img, Inches(7.0), Inches((7.5 - vid_h) / 2),
                      width=Inches(5.88), height=Inches(vid_h))

logo_small(s3)
eyebrow(s3, 'Technical Approach')

txt(s3, 'A CV pipeline that reads a\nconstruction site like a 3D scene.',
    Inches(0.65), Inches(0.78), Inches(6.0), Inches(1.28),
    sz=27, bold=True, color=TEXT)

hrule(s3, Inches(2.2), w=Inches(5.95))

# Screenshot label
txt(s3, 'Landex annotation UI — component verification against approved submittal',
    Inches(7.0), Inches(5.98), Inches(5.88), Inches(0.35),
    sz=8, color=TEXT_SEC, italic=True, align=PP_ALIGN.CENTER)

steps = [
    ('01  Video Ingestion',
     'Walk-video of the recently installed area captured on a phone. '
     'No special rig, no LiDAR, no structured-light sensor.'),
    ('02  Frame Analysis & Detection',
     'Model identifies installed MEP and structural components across frames, '
     'reconciling detections against the IFC/Revit element list from the approved submittal.'),
    ('03  Human-in-the-Loop Confirmation',
     'Engineer reviews auto-tagged results, corrects misses or false positives. '
     'Corrections feed back into the model.'),
    ('04  BIM Write-Back',
     'Each verified element stamped into the 3D model with timestamp, video reference, '
     'and submitter ID — queryable at the element level.'),
]

y = Inches(2.38)
for title, body in steps:
    txt(s3, title, Inches(0.65), y, Inches(5.92), Inches(0.3),
        sz=10.5, bold=True, color=GREEN)
    txt(s3, body, Inches(0.65), y + Inches(0.27), Inches(5.92), Inches(0.68),
        sz=9.5, color=TEXT_SEC)
    y += Inches(1.12)

txt(s3,
    'Core challenge: multi-view component recognition under occlusion, '
    'across varied lighting, without point-cloud input.',
    Inches(0.65), Inches(6.96), Inches(5.92), Inches(0.36),
    sz=8, color=TEXT_SEC, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — THE MARKET
# ═══════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
fill_rect(s4, 0, 0, W, H)
logo_small(s4)
eyebrow(s4, 'The Market')

txt(s4, 'Bad construction data is a\nknown, quantified problem.',
    Inches(0.65), Inches(0.78), Inches(12), Inches(1.2),
    sz=30, bold=True, color=TEXT)

hrule(s4, Inches(2.18))

stats = [
    ('$1.85T',  'Lost to bad construction data globally in 2020',                            'Autodesk / FMI'),
    ('98%',     'Of megaprojects run over 30% on cost',                                       'McKinsey'),
    ('5–6%',    'Typical contractor pre-tax margin — a few points of rework erases the job',  ''),
    ('3–7%',    'Of project cost lost to billing and progress errors caught late',             ''),
]

cw  = Inches(5.9)
ch  = Inches(2.12)
gap = Inches(0.23)
xs  = [Inches(0.65), Inches(0.65) + cw + gap]
ys  = [Inches(2.38), Inches(2.38) + ch + gap]

for (yi, yv) in enumerate(ys):
    for (xi, xv) in enumerate(xs):
        fig, label, src = stats[yi * 2 + xi]
        fill_rect(s4, xv, yv, cw, ch, color=SURFACE, border=BORDER)
        txt(s4, fig, xv + Inches(0.25), yv + Inches(0.17), cw - Inches(0.3), Inches(0.72),
            sz=40, bold=True, color=GREEN)
        txt(s4, label, xv + Inches(0.25), yv + Inches(0.95), cw - Inches(0.3), Inches(0.68),
            sz=10, color=TEXT_SEC)
        if src:
            txt(s4, src, xv + Inches(0.25), yv + Inches(1.77), cw - Inches(0.3), Inches(0.28),
                sz=8, color=TEXT_SEC, italic=True)

txt(s4,
    'The verification step exists on every job. The software to close the loop does not.',
    Inches(0.65), Inches(7.07), Inches(12), Inches(0.33),
    sz=9, color=TEXT_SEC, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TEAM
# ═══════════════════════════════════════════════════════════════════════════════
s5 = new_slide()
fill_rect(s5, 0, 0, W, H)
logo_small(s5)
eyebrow(s5, 'Who We Are')

txt(s5, 'MIT and Stanford.\nRobotics, ML, and CV for the physical world.',
    Inches(0.65), Inches(0.72), Inches(11.5), Inches(1.2),
    sz=28, bold=True, color=TEXT)

hrule(s5, Inches(2.1))

members = [
    ('Allen Chen', 'Co-founder',
     A + 'team/allen.png',
     'MIT Mechanical Engineering \'26. Robotics focus. Brings domain understanding '
     'of how physical systems are specified and built — the gap most CV teams miss on '
     'construction problems.'),
    ('Auddithio Nag', 'Co-founder',
     A + 'team/auddi.png',
     'Stanford CS MS \'26. ML research in geospatial analysis, satellite flood prediction, '
     'medical imaging, and 3D scene understanding. Spatial inference from sensor data is '
     'the foundation of Landex\'s CV pipeline.'),
    ('Beckett Devoe', 'Co-founder',
     A + 'team/beckett.png',
     'MIT AI & Decision Making \'26. Computer vision for the physical world: shellfish '
     'health at MIT Sea Grant, ocean sensing in Norway. Pointing the same tools at '
     'the built environment.'),
]

cw  = Inches(3.92)
gap = Inches(0.22)
xs  = [Inches(0.65), Inches(0.65) + cw + gap, Inches(0.65) + 2 * (cw + gap)]

for x, (name, role, photo_path, bio) in zip(xs, members):
    # Card
    fill_rect(s5, x, Inches(2.28), cw, Inches(4.95), color=SURFACE, border=BORDER)

    # Headshot — square crop, centered in card
    ph_w = Inches(1.95)
    ph_x = x + (cw - ph_w) / 2
    sq = sq_crop(photo_path)
    s5.shapes.add_picture(sq, ph_x, Inches(2.45), width=ph_w, height=ph_w)

    txt(s5, name, x + Inches(0.18), Inches(4.53),
        cw - Inches(0.22), Inches(0.38), sz=13, bold=True, color=TEXT)
    txt(s5, role, x + Inches(0.18), Inches(4.91),
        cw - Inches(0.22), Inches(0.28), sz=9, color=GREEN)
    txt(s5, bio, x + Inches(0.18), Inches(5.26),
        cw - Inches(0.22), Inches(1.82), sz=9.5, color=TEXT_SEC)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — THE ASK
# ═══════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
fill_rect(s6, 0, 0, W, H)

# Right panel — SURFACE + dashboard screenshot
fill_rect(s6, Inches(6.85), 0, Inches(6.48), H, color=SURFACE, border=BORDER)

dash_h = 5.88 / (2574 / 1820)  # = 4.162"
dash_img = prep_img(A + 'product-dashboard.png', 5.88, dash_h)
s6.shapes.add_picture(dash_img, Inches(7.0), Inches((7.5 - dash_h) / 2),
                      width=Inches(5.88), height=Inches(dash_h))

logo_small(s6)
eyebrow(s6, 'The Ask')

txt(s6, 'One scope, one pay cycle.\nProve it where the money is moving.',
    Inches(0.65), Inches(0.78), Inches(6.0), Inches(1.28),
    sz=27, bold=True, color=TEXT)

hrule(s6, Inches(2.2), w=Inches(5.95))

# Dashboard label
txt(s6, 'Landex director view — verified BIM with element-level install records',
    Inches(7.0), Inches(6.1), Inches(5.88), Inches(0.35),
    sz=8, color=TEXT_SEC, italic=True, align=PP_ALIGN.CENTER)

asks = [
    ('Pilot Structure',
     'One scope on a live job. Fixed fee, fixed window. Video in, verified model out. '
     'No new field hardware, no open-ended integration.'),
    ('Why NVIDIA',
     'Our inference workload is GPU-bound: multi-frame component detection under occlusion, '
     'real-time reconciliation against large IFC element sets. NVIDIA Insight gives us '
     'compute credits and the CV/spatial computing stack to push this to production scale.'),
    ('What We\'re Asking For',
     'GPU compute credits, access to NVIDIA\'s CV and spatial computing tools, and technical '
     'mentorship as we move from pilot to fleet deployment across active jobsites.'),
]

y = Inches(2.42)
for label, body in asks:
    green_tick(s6, Inches(0.65), y + Inches(0.09))
    txt(s6, label, Inches(0.9), y, Inches(5.65), Inches(0.32),
        sz=12, bold=True, color=SAND)
    txt(s6, body, Inches(0.9), y + Inches(0.31), Inches(5.65), Inches(0.78),
        sz=10.5, color=TEXT_SEC)
    y += Inches(1.3)

hrule(s6, Inches(6.75), w=Inches(5.95))
txt(s6, 'allen@landexsystems.com  ·  landexsystems.com',
    Inches(0.65), Inches(6.88), Inches(6), Inches(0.38),
    sz=10, color=TEXT_SEC)


# ── SAVE ─────────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f'Saved: {OUT}')
